#!/usr/bin/env python3
"""Verify the proposal content"""
from pptx import Presentation

PATH = "/Users/sheikhown/Library/CloudStorage/OneDrive-Personal/Rank Ray - Company/RR - Office Works/Proposals/SEO/SEO Proposal - MySportsInjury - Hamza.pptx"
prs = Presentation(PATH)

print(f"Total slides: {len(prs.slides)}\n")
for i, slide in enumerate(prs.slides):
    print(f"=== SLIDE {i+1} ===")
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                # Truncate long text for display
                if len(text) > 300:
                    print(f"  {text[:300]}...")
                else:
                    print(f"  {text}")
    print()
