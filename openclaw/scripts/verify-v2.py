#!/usr/bin/env python3
"""Verify the new proposal structure"""
from pptx import Presentation

PATH = "/Users/sheikhown/Library/CloudStorage/OneDrive-Personal/Rank Ray - Company/RR - Office Works/Proposals/SEO/SEO Proposal - MySportsInjury - Hamza.pptx"
try:
    prs = Presentation(PATH)
    print(f"✅ File opens successfully")
    print(f"Total slides: {len(prs.slides)}\n")
    
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text and len(text) > 2:
                    texts.append(text[:80])
        
        preview = " | ".join(texts[:3]) if texts else "[images/design only]"
        print(f"Slide {i+1}: {preview}")
except Exception as e:
    print(f"❌ Error: {e}")
