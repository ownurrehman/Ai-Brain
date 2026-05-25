#!/usr/bin/env python3
"""
Create SEO Proposal for mysportsinjury.co.uk
Based on Rank Ray SEO Profile template + SEMrush data
"""

import shutil
from pptx import Presentation
from datetime import datetime

# Copy template
TEMPLATE_PATH = "/Users/sheikhown/Library/CloudStorage/OneDrive-Personal/Rank Ray - Company/RR - Office Works/Proposals/SEO/Rank Ray SEO Profile.pptx"
OUTPUT_PATH = "/Users/sheikhown/Library/CloudStorage/OneDrive-Personal/Rank Ray - Company/RR - Office Works/Proposals/SEO/SEO Proposal - MySportsInjury - Hamza.pptx"

shutil.copy(TEMPLATE_PATH, OUTPUT_PATH)

prs = Presentation(OUTPUT_PATH)

# Content for each slide
slides_content = [
    {
        "idx": 0,
        "title": "SEO PROPOSAL",
        "subtitle": "For mysportsinjury.co.uk\nPrepared by Rank Ray\nMay 2026"
    },
    {
        "idx": 1,
        "title": "ABOUT THE CLIENT",
        "body": "Website: mysportsinjury.co.uk\nIndustry: Sports Physiotherapy & Rehabilitation\nLocation: Manchester, UK\nContact: Hamza (Representative)\n\nMy Sports Injury is a leading sports physiotherapy clinic based in Manchester, UK. They specialize in treating sports injuries, offering physiotherapy, sports massage, rehabilitation programs, and recovery products. Their target audience includes athletes, fitness enthusiasts, and individuals seeking professional physiotherapy care in the UK."
    },
    {
        "idx": 2,
        "title": "CURRENT SEO PERFORMANCE",
        "body": "Based on SEMrush Analysis (UK Database):\n\n• Domain Rank: #244,526\n• Organic Keywords: 437 keywords\n• Organic Traffic: ~1,339 visits/month\n• Traffic Value: £545/month\n• Paid Ads: None running\n\nCurrent Keyword Positions:\n• 'sports physio manchester' - Position #1 (110 vol)\n• 'medial tibial periostitis' - Position #3 (480 vol)\n• 'lumbar lordosis' - Position #11 (4,400 vol)\n• 'thoracic back pain' - Position #17 (2,400 vol)\n• 'ankle syndesmosis' - Position #8 (590 vol)\n\nKey Finding: Strong content but significant room for ranking improvements across high-volume keywords."
    },
    {
        "idx": 3,
        "title": "SEO OPPORTUNITIES & STRATEGY",
        "body": "Our 3-Month SEO Plan:\n\n1. TECHNICAL SEO AUDIT\n   • Fix crawl errors & site speed\n   • Mobile optimization\n   • Schema markup for local business\n\n2. ON-PAGE OPTIMIZATION\n   • Optimize existing 437+ keyword targets\n   • Improve title tags & meta descriptions\n   • Internal linking structure\n\n3. CONTENT STRATEGY\n   • Target high-volume terms (lumbar lordosis 4,400/mo)\n   • Create condition-specific landing pages\n   • Sports injury guides & blog content\n\n4. LOCAL SEO\n   • Google Business Profile optimization\n   • Local citations for Manchester physio\n   • Location-based landing pages\n\n5. BACKLINK BUILDING\n   • Sports & health industry outreach\n   • Guest posting on fitness/physio blogs\n   • Local directory submissions"
    },
    {
        "idx": 4,
        "title": "PRICING & CONTRACT TERMS",
        "body": "INITIAL 3-MONTH CAMPAIGN\n\nMonthly Investment: £500 GBP\nDuration: 3 Months\nTotal Investment: £1,500 GBP\n\nWhat's Included:\n• Full technical SEO audit & fixes\n• On-page optimization for all key pages\n• 8+ new SEO-optimized articles/month\n• Local SEO & GBP management\n• Monthly backlink acquisition (5-10 quality links)\n• Weekly rank tracking & monthly reporting\n\nUPGRADE PATH:\nIf Rank Ray delivers measurable results within 3 months,\nthe contract upgrades to a 1-Year Agreement\nwith adjusted pricing based on performance.\n\nPayment: Monthly in advance\nStart Date: Upon agreement signature"
    },
    {
        "idx": 5,
        "title": "EXPECTED RESULTS (3 MONTHS)",
        "body": "Conservative Projections for mysportsinjury.co.uk:\n\n• Organic Keywords: 437 → 700+ (60% increase)\n• Organic Traffic: 1,339 → 3,000+ visits/month\n• Top 3 Rankings: 5 → 25+ keywords\n• Top 10 Rankings: 15 → 50+ keywords\n\nPriority Keywords to Improve:\n• 'lumbar lordosis' (4,400/mo) - Target: Top 5\n• 'thoracic back pain' (2,400/mo) - Target: Top 10\n• 'sports physiotherapy manchester' - Target: #1\n• 'ankle syndesmosis' (590/mo) - Target: Top 3\n\nSuccess Metrics:\n• Increased appointment bookings\n• Higher local search visibility\n• Improved domain authority\n• Measurable ROI on organic traffic"
    },
    {
        "idx": 6,
        "title": "WHY RANK RAY?",
        "body": "Your Dedicated SEO Partner:\n\n• Specialized in local & healthcare SEO\n• Proven track record with UK businesses\n• Transparent reporting & communication\n• Data-driven approach using SEMrush, Ahrefs\n• Focus on ROI and measurable results\n\nOur Process:\n1. Audit & Strategy (Week 1-2)\n2. Quick Wins & Technical Fixes (Week 3-4)\n3. Content & Link Building (Month 2-3)\n4. Reporting & Optimization (Ongoing)\n\nContact:\nOwn-ur-Rehman Sheikh, CEO\nRank Ray Digital Agency\nEmail: contact@rankray.com\nWeb: rankray.com"
    },
]

def get_placeholder_type(shape):
    """Safely get placeholder type or None"""
    try:
        if shape.is_placeholder:
            return shape.placeholder_format.type
    except:
        pass
    return None

# Update slides
for slide_info in slides_content:
    idx = slide_info["idx"]
    if idx < len(prs.slides):
        slide = prs.slides[idx]
        shapes = slide.shapes
        
        # Find title and body shapes
        title_shape = None
        body_shape = None
        
        for shape in shapes:
            if shape.has_text_frame:
                ph_type = get_placeholder_type(shape)
                if ph_type == 1 or ph_type == 3:  # TITLE or CENTER_TITLE
                    title_shape = shape
                elif ph_type == 2:  # BODY
                    body_shape = shape
        
        # If no placeholders found, find first shape with text_frame as title, second as body
        if title_shape is None:
            for shape in shapes:
                if shape.has_text_frame:
                    title_shape = shape
                    break
        if body_shape is None:
            found_title = False
            for shape in shapes:
                if shape.has_text_frame:
                    if not found_title:
                        found_title = True
                        continue
                    body_shape = shape
                    break
        
        # Update content
        if title_shape and "title" in slide_info:
            title_shape.text_frame.text = slide_info["title"]
        if body_shape:
            if "subtitle" in slide_info:
                body_shape.text_frame.text = slide_info["subtitle"]
            elif "body" in slide_info:
                body_shape.text_frame.text = slide_info["body"]

prs.save(OUTPUT_PATH)
print(f"✅ Proposal created: {OUTPUT_PATH}")
print(f"   Total slides: {len(prs.slides)}")
