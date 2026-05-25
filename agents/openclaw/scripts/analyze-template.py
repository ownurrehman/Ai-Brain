#!/usr/bin/env python3
"""Deep analysis of the PPTX template structure"""
from pptx import Presentation
from pptx.util import Inches, Pt

TEMPLATE_PATH = "/Users/sheikhown/Library/CloudStorage/OneDrive-Personal/Rank Ray - Company/RR - Office Works/Proposals/SEO/Rank Ray SEO Profile.pptx"
prs = Presentation(TEMPLATE_PATH)

print(f"Slide width: {prs.slide_width.inches:.2f} inches")
print(f"Slide height: {prs.slide_height.inches:.2f} inches")
print(f"Total slides: {len(prs.slides)}")
print(f"Total layouts: {len(prs.slide_layouts)}")
print()

# List all layouts
print("=== AVAILABLE LAYOUTS ===")
for i, layout in enumerate(prs.slide_layouts):
    placeholders = list(layout.placeholders) if layout.placeholders else []
    ph_info = [f"{p.placeholder_format.type}({p.placeholder_format.idx})" for p in placeholders]
    print(f"Layout {i}: '{layout.name}' - {len(placeholders)} placeholders: {', '.join(ph_info)}")
print()

# Analyze each slide in detail
for i, slide in enumerate(prs.slides):
    print(f"\n{'='*60}")
    print(f"SLIDE {i+1}")
    print(f"Layout: {slide.slide_layout.name}")
    print(f"Shapes count: {len(slide.shapes)}")
    print()
    
    for j, shape in enumerate(slide.shapes):
        shape_type = str(shape.shape_type).replace("MSO_SHAPE_TYPE.", "")
        
        # Check if placeholder
        is_ph = False
        ph_type = "N/A"
        ph_idx = "N/A"
        try:
            if shape.is_placeholder:
                is_ph = True
                ph_type = shape.placeholder_format.type
                ph_idx = shape.placeholder_format.idx
        except:
            pass
        
        # Get text preview
        text_preview = ""
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                text_preview = text[:100].replace("\n", "\\n")
        
        # Get position/size
        pos = f"left={shape.left.inches:.2f}, top={shape.top.inches:.2f}, w={shape.width.inches:.2f}, h={shape.height.inches:.2f}"
        
        print(f"  Shape {j}: type={shape_type}, placeholder={is_ph}, ph_type={ph_type}, ph_idx={ph_idx}")
        print(f"           pos=[{pos}]")
        if text_preview:
            print(f"           text=\"{text_preview}\"")
    print()
