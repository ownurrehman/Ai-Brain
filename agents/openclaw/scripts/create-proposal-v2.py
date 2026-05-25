#!/usr/bin/env python3
"""
Create SEO Proposal for mysportsinjury.co.uk
PRESERVING template structure, design, colors, and portfolio slides
"""

import shutil
from pptx import Presentation
from pptx.util import Inches, Pt

TEMPLATE_PATH = "/Users/sheikhown/Library/CloudStorage/OneDrive-Personal/Rank Ray - Company/RR - Office Works/Proposals/SEO/Rank Ray SEO Profile.pptx"
OUTPUT_PATH = "/Users/sheikhown/Library/CloudStorage/OneDrive-Personal/Rank Ray - Company/RR - Office Works/Proposals/SEO/SEO Proposal - MySportsInjury - Hamza.pptx"

shutil.copy(TEMPLATE_PATH, OUTPUT_PATH)
prs = Presentation(OUTPUT_PATH)

# Layout mapping from template
LAYOUT_TITLE = prs.slide_layouts[0]           # TITLE: CENTER_TITLE + SUBTITLE
LAYOUT_TITLE_BODY = prs.slide_layouts[1]        # TITLE_AND_BODY: BODY + TITLE
LAYOUT_SECTION_HDR = prs.slide_layouts[2]       # SECTION_HEADER: TITLE + SUBTITLE
LAYOUT_ONE_COL = prs.slide_layouts[3]         # ONE_COLUMN_TEXT: TITLE + BODY
LAYOUT_TITLE_ONLY = prs.slide_layouts[4]      # TITLE_ONLY: just TITLE
LAYOUT_BIG_NUM = prs.slide_layouts[8]         # BIG_NUMBER: TITLE + BODY
LAYOUT_TWO_COL = prs.slide_layouts[13]        # TITLE_AND_TWO_COLUMNS

def set_text_in_placeholder(slide, placeholder_type, text):
    """Find placeholder by type and set text"""
    for shape in slide.shapes:
        if shape.is_placeholder:
            if shape.placeholder_format.type == placeholder_type:
                shape.text_frame.text = text
                return True
    return False

def add_title_slide(prs, layout, title, subtitle):
    """Add a title slide with proper text"""
    slide = prs.slides.add_slide(layout)
    set_text_in_placeholder(slide, 3, title)  # CENTER_TITLE
    set_text_in_placeholder(slide, 4, subtitle)  # SUBTITLE
    return slide

def add_content_slide(prs, layout, title, body):
    """Add a content slide with title and body"""
    slide = prs.slides.add_slide(layout)
    set_text_in_placeholder(slide, 1, title)  # TITLE
    set_text_in_placeholder(slide, 2, body)   # BODY
    return slide

def add_section_slide(prs, layout, title, subtitle):
    """Add a section header slide"""
    slide = prs.slides.add_slide(layout)
    set_text_in_placeholder(slide, 1, title)  # TITLE
    set_text_in_placeholder(slide, 4, subtitle)  # SUBTITLE
    return slide

# ============================================================
# STRATEGY: Insert new proposal slides at position 1 (after title)
# Keep original slides 2-14 intact (portfolio, packages, contact, thank you)
# ============================================================

# We need to insert slides at specific positions.
# python-pptx doesn't support insert, so we'll rebuild the deck.

# Collect all original slides (we'll keep most of them)
original_slides_data = []
for slide in prs.slides:
    original_slides_data.append(slide)

# Now create a NEW presentation from the same template
import os
os.remove(OUTPUT_PATH)
shutil.copy(TEMPLATE_PATH, OUTPUT_PATH)

# Re-open fresh
prs2 = Presentation(TEMPLATE_PATH)

# We'll rebuild slide by slide
# First, let me extract the XML of slides to copy them properly

# Actually, the simplest approach: create new slides and move them
# But python-pptx can't easily reorder slides.

# Alternative: Create a new presentation from scratch using the template's layouts
# and copy only the slides we want, then add new ones in between

# Let me try a different approach: use the original template, 
# add new slides at the end, then reorder using XML manipulation

# Actually, simplest: work with the existing deck, add new slides,
# then use XML to move them to the right positions.

from pptx.opc.constants import RELATIONSHIP_TYPE as RT

prs = Presentation(OUTPUT_PATH)

# Add all our proposal-specific slides at the end first
proposal_slides = []

# Slide A: SEO Proposal Title (for mysportsinjury)
s1 = prs.slides.add_slide(LAYOUT_TITLE)
set_text_in_placeholder(s1, 3, "SEO PROPOSAL")
set_text_in_placeholder(s1, 4, "For mysportsinjury.co.uk\nManchester, UK\nMay 2026")
proposal_slides.append(s1)

# Slide B: About the Client (using ONE_COLUMN_TEXT layout)
s2 = prs.slides.add_slide(LAYOUT_ONE_COL)
set_text_in_placeholder(s2, 1, "ABOUT THE CLIENT")
body_text = """Website: mysportsinjury.co.uk
Industry: Sports Physiotherapy & Rehabilitation
Location: Manchester, UK
Contact: Hamza (Representative)

My Sports Injury is a leading sports physiotherapy clinic in Manchester, UK. They treat sports injuries, offer physiotherapy, sports massage, rehabilitation programs, and recovery products for athletes and fitness enthusiasts."""
set_text_in_placeholder(s2, 2, body_text)
proposal_slides.append(s2)

# Slide C: Current SEO Performance (BIG_NUMBER layout for stats)
s3 = prs.slides.add_slide(LAYOUT_BIG_NUM)
set_text_in_placeholder(s3, 1, "CURRENT SEO PERFORMANCE")
perf_text = """Based on SEMrush Analysis (UK Database):

• Domain Rank: #244,526
• Organic Keywords: 437 keywords
• Organic Traffic: ~1,339 visits/month
• Traffic Value: £545/month

Key Rankings:
• 'sports physio manchester' — Position #1
• 'medial tibial periostitis' — Position #3 (480/mo vol)
• 'lumbar lordosis' — Position #11 (4,400/mo vol)
• 'thoracic back pain' — Position #17 (2,400/mo vol)"""
set_text_in_placeholder(s3, 2, perf_text)
proposal_slides.append(s3)

# Slide D: SEO Opportunities & Strategy (TITLE_AND_BODY)
s4 = prs.slides.add_slide(LAYOUT_TITLE_BODY)
set_text_in_placeholder(s4, 1, "SEO OPPORTUNITIES & STRATEGY")
strategy_text = """Our 3-Month SEO Plan:

1. TECHNICAL SEO AUDIT
   • Fix crawl errors & improve site speed
   • Mobile optimization & Core Web Vitals
   • Schema markup for LocalBusiness

2. ON-PAGE OPTIMIZATION
   • Optimize 437+ existing keyword targets
   • Improve title tags, meta descriptions
   • Strengthen internal linking structure

3. CONTENT STRATEGY
   • Target high-volume terms (lumbar lordosis 4,400/mo)
   • Create condition-specific landing pages
   • Publish sports injury guides & rehab content

4. LOCAL SEO
   • Google Business Profile optimization
   • Manchester physiotherapy citations
   • Location-based landing pages

5. BACKLINK BUILDING
   • Sports & health industry outreach
   • Guest posting on fitness/physio blogs
   • Local directory submissions"""
set_text_in_placeholder(s4, 2, strategy_text)
proposal_slides.append(s4)

# Slide E: Pricing & Contract Terms (SECTION_HEADER layout)
s5 = prs.slides.add_slide(LAYOUT_SECTION_HDR)
set_text_in_placeholder(s5, 1, "PRICING & CONTRACT TERMS")
pricing_text = """Initial 3-Month Campaign

Monthly Investment: £500 GBP
Duration: 3 Months
Total: £1,500 GBP

Includes:
• Full technical SEO audit & fixes
• On-page optimization for key pages
• 8+ new SEO-optimized articles/month
• Local SEO & GBP management
• Monthly backlink acquisition (5-10 quality links)
• Weekly rank tracking & monthly reporting

UPGRADE PATH:
If Rank Ray delivers measurable results within 3 months, the contract upgrades to a 1-Year Agreement with adjusted pricing based on performance.

Payment: Monthly in advance | Start: Upon signature"""
set_text_in_placeholder(s5, 4, pricing_text)
proposal_slides.append(s5)

# Slide F: Expected Results (TITLE_AND_BODY)
s6 = prs.slides.add_slide(LAYOUT_TITLE_BODY)
set_text_in_placeholder(s6, 1, "EXPECTED RESULTS (3 MONTHS)")
results_text = """Conservative Projections:

• Organic Keywords: 437 → 700+ (60% increase)
• Organic Traffic: 1,339 → 3,000+ visits/month
• Top 3 Rankings: 5 → 25+ keywords
• Top 10 Rankings: 15 → 50+ keywords

Priority Keywords to Target:
• 'lumbar lordosis' (4,400/mo) — Target: Top 5
• 'thoracic back pain' (2,400/mo) — Target: Top 10
• 'sports physiotherapy manchester' — Target: #1
• 'ankle syndesmosis' (590/mo) — Target: Top 3

Success Metrics:
• Increased appointment bookings
• Higher local search visibility
• Improved domain authority
• Measurable ROI on organic traffic"""
set_text_in_placeholder(s6, 2, results_text)
proposal_slides.append(s6)

# Now we have:
# Slides 1-14: Original template (title, about, vision, experience, why us, portfolio x3, SEO text, how we do it, purpose, packages, contact, thank you)
# Slides 15-20: New proposal slides

# We want the order:
# 1. SEO Proposal Title (new)
# 2. About the Client (new)
# 3. Current SEO Performance (new)
# 4. SEO Opportunities (new)
# 5. Pricing (new)
# 6. Expected Results (new)
# 7. Original title slide (Rank Ray intro)
# 8. Original about/vision/experience/why us
# 9. Portfolio slides (keep screenshots)
# 10. SEO services slide (keep)
# 11. How we do it (keep)
# 12. Purpose (keep)
# 13. Packages (update pricing for this client?)
# 14. Contact (update for this proposal?)
# 15. Thank You (keep)

# The XML approach to reorder slides:
# Get the slide id list from presentation.xml

# Get slide IDs in current order
slide_id_list = prs.part._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}sldIdLst')

# Current order: original slides 1-14, then new slides 15-20
# Desired order: new slides first (positions 15-20 → 1-6), then original slides 1-14

# Get all slide IDs
slide_ids = list(slide_id_list)
print(f"Total slides: {len(slide_ids)}")

# New slides are at indices 14-19 (0-based), original at 0-13
new_slide_ids = slide_ids[14:20]  # 6 new proposal slides
orig_slide_ids = slide_ids[0:14]  # 14 original template slides

# Clear current list
for child in list(slide_id_list):
    slide_id_list.remove(child)

# Re-add in desired order: new slides first, then originals
for s in new_slide_ids:
    slide_id_list.append(s)
for s in orig_slide_ids:
    slide_id_list.append(s)

# Save
prs.save(OUTPUT_PATH)
print(f"✅ Proposal created: {OUTPUT_PATH}")
print(f"   Total slides: {len(prs.slides)}")
print(f"   Order: 6 proposal slides first, then 14 original template slides")
